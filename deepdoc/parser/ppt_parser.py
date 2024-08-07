from io import BytesIO
from pptx import Presentation


class RAGFlowPptParser(object):
    """
    RAG流程PPT解析器类，用于解析特定格式的PPT，提取其中的信息。
    """

    def __init__(self):
        """
        初始化解析器。
        """
        super().__init__()

    def __extract(self, shape):
        """
        内部方法，用于递归地从一个形状中提取文本信息。

        参数:
        - shape: pptx库中的Shape对象，表示PPT中的一个形状，可以是文本框、表格等。

        返回:
        - str: 从给定形状中提取的文本信息，格式化为适合进一步处理的字符串。
        """
        # 检查形状是否为表格类型
        if shape.shape_type == 19:
            tb = shape.table
            rows = []
            # 遍历表格行，忽略表头，拼接每行的关键信息
            for i in range(1, len(tb.rows)):
                rows.append("; ".join([tb.cell(0, j).text + ": " + tb.cell(i, j).text
                                       for j in range(len(tb.columns)) if tb.cell(i, j)]))
            return "\n".join(rows)

        # 检查形状是否包含文本框
        if shape.has_text_frame:
            return shape.text_frame.text

        # 检查形状是否为组合形状，递归处理其中的每个子形状
        if shape.shape_type == 6:
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                t = self.__extract(p)
                if t:
                    texts.append(t)
            return "\n".join(texts)

    def __call__(self, fnm, from_page, to_page, callback=None):
        """
        自定义调用方法，用于提取指定 PowerPoint 文件中指定页码范围内的文本内容。

        参数:
        - fnm: PowerPoint 文件名或文件对象。如果是字符串，则表示文件路径；如果是 BytesIO 对象，则表示文件内容。
        - from_page: 起始提取的页码。
        - to_page: 结束提取的页码。
        - callback: 回调函数，用于处理每一页提取到的文本。

        返回:
        - 一个包含指定页码范围内所有页面提取到的文本内容的列表，每个元素是一个字符串，代表一页的内容。
        """
        # 根据 fnm 的类型，创建 Presentation 对象，可以处理文件路径或 BytesIO 对象
        ppt = Presentation(fnm) if isinstance(
            fnm, str) else Presentation(
            BytesIO(fnm))
        # 初始化用于存储所有页面文本内容的列表
        txts = []
        # 记录演示文稿的总页数
        self.total_page = len(ppt.slides)
        # 遍历每一页，提取文本内容
        for i, slide in enumerate(ppt.slides):
            # 跳过起始页码之前的内容
            if i < from_page:
                continue
            # 结束提取页码之后的内容
            if i >= to_page:
                break
            # 初始化当前页面文本内容的列表
            texts = []
            # 对当前页面的所有形状进行排序，按其在页面上的位置（顶部和左侧）
            for shape in sorted(
                    slide.shapes, key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left)):
                txt = self.__extract(shape)
                # 如果提取到文本，则添加到当前页面的文本内容列表中
                if txt:
                    texts.append(txt)
            # 将当前页面的文本内容列表合并成一个字符串，加入到所有页面文本内容的列表中
            txts.append("\n".join(texts))

        return txts
